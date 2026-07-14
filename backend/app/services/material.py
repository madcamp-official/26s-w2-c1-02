"""PDF·PPTX → slides.json 파서 (PyMuPDF · python-pptx).

발표 자료(PDF/PPTX)에서 페이지별 텍스트를 추출해 materials.slides JSONB 형식
(db-schema §6.1: [{"page": 1, "text": "..."}]) 으로 변환한다.

팀원2의 백그라운드 파싱 잡(POST /sessions/{id}/material → 202)에서 호출:

    try:
        slides = parse_material_to_slides(data, ext)  # 확장자로 파서 선택
    except UnprocessableMaterialError:
        # → status="failed", error_code="UNPROCESSABLE_PDF|PPTX" (api-spec §4.2)
    except MaterialParseError:
        # → 재시도 또는 실패 처리

동기(CPU-bound) 함수이므로 async 잡에서는 run_in_executor 등으로 감싼다.
"""

import io
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# api-spec §1.3: 자료 최대 50페이지/슬라이드 (업로드 시 FE/라우터가 1차 검증)
MAX_PAGES = 50


class MaterialParseError(Exception):
    """자료 파싱 실패(손상 파일 등). 잡에서 실패/재시도 처리."""


class UnprocessableMaterialError(MaterialParseError):
    """텍스트 추출 불가(스캔본·이미지 전용·암호화). → 422 UNPROCESSABLE_*."""


class PdfParseError(MaterialParseError):
    """PDF 파싱 실패(손상 파일 등). 잡에서 실패/재시도 처리."""


class UnprocessablePdfError(PdfParseError, UnprocessableMaterialError):
    """텍스트 추출 불가(스캔본·이미지 PDF·암호화). → 422 UNPROCESSABLE_PDF."""


class PptxParseError(MaterialParseError):
    """PPTX 파싱 실패(손상 파일·PPTX 아님 등). 잡에서 실패/재시도 처리."""


class UnprocessablePptxError(PptxParseError, UnprocessableMaterialError):
    """텍스트 추출 불가(이미지 전용 덱·슬라이드 수 초과). → 422 UNPROCESSABLE_PPTX."""


def parse_material_to_slides(data: bytes, ext: str) -> list[dict]:
    """확장자에 맞는 파서로 위임한다 — 파싱 잡의 단일 진입점.

    Args:
        data: 자료 파일 bytes.
        ext: 파일 확장자 ("pdf" | "pptx", 대소문자·선행 점 무관).
    """
    ext = ext.lower().lstrip(".")
    if ext == "pdf":
        return parse_pdf_to_slides(data)
    if ext == "pptx":
        return parse_pptx_to_slides(data)
    raise MaterialParseError(f"지원하지 않는 자료 형식: {ext}")


def parse_pdf_to_slides(pdf: bytes | str | Path) -> list[dict]:
    """PDF에서 페이지별 텍스트를 추출한다.

    Args:
        pdf: PDF 파일의 bytes 또는 경로.

    Returns:
        [{"page": 1, "text": "..."}] — page는 1부터, 텍스트 없는 페이지는
        text=""로 포함해 page 번호가 실제 PDF와 항상 일치하게 유지한다
        (질문 evidence.slides가 이 번호를 참조, db-schema §6.3).

    Raises:
        UnprocessablePdfError: 전 페이지에서 텍스트를 못 얻은 경우(스캔본),
            암호화된 PDF, 페이지 수 초과.
        PdfParseError: 파일이 PDF가 아니거나 손상된 경우.
    """
    try:
        if isinstance(pdf, bytes):
            doc = fitz.open(stream=pdf, filetype="pdf")
        else:
            doc = fitz.open(str(pdf))
    except Exception as e:
        raise PdfParseError(f"PDF를 열 수 없음: {e}") from e

    try:
        if doc.needs_pass:
            raise UnprocessablePdfError("암호화된 PDF")
        if doc.page_count > MAX_PAGES:
            raise UnprocessablePdfError(f"페이지 수 초과: {doc.page_count} > {MAX_PAGES}")

        slides = [
            {"page": i + 1, "text": _clean_text(page.get_text("text"))}
            for i, page in enumerate(doc)
        ]
    except PdfParseError:
        raise
    except Exception as e:
        # fitz.open이 성공해도 특정 페이지의 스트림이 손상돼 있으면 get_text가
        # raw RuntimeError/FileDataError를 던진다 — 문서화된 계약으로 변환.
        raise PdfParseError(f"페이지 텍스트 추출 실패: {e}") from e
    finally:
        doc.close()

    if not slides or all(not s["text"] for s in slides):
        raise UnprocessablePdfError("텍스트 레이어 없음(스캔본 추정)")
    return slides


def parse_pptx_to_slides(pptx: bytes | str | Path) -> list[dict]:
    """PPTX에서 슬라이드별 텍스트를 추출한다.

    Args:
        pptx: PPTX 파일의 bytes 또는 경로.

    Returns:
        [{"page": 1, "text": "..."}] — parse_pdf_to_slides와 동일 계약.
        텍스트 없는 슬라이드도 text=""로 포함해 page 번호가 실제 덱과 항상
        일치하게 유지한다 (질문 evidence.slides가 이 번호를 참조).

    Raises:
        UnprocessablePptxError: 전 슬라이드에서 텍스트를 못 얻은 경우
            (이미지 전용 덱), 슬라이드 수 초과.
        PptxParseError: 파일이 PPTX가 아니거나 손상된 경우
            (암호화된 PPTX는 zip 컨테이너가 아니라서 열기 단계에서 여기로 온다).
    """
    try:
        if isinstance(pptx, bytes):
            prs = Presentation(io.BytesIO(pptx))
        else:
            prs = Presentation(str(pptx))
    except Exception as e:
        raise PptxParseError(f"PPTX를 열 수 없음: {e}") from e

    if len(prs.slides) > MAX_PAGES:
        raise UnprocessablePptxError(f"슬라이드 수 초과: {len(prs.slides)} > {MAX_PAGES}")

    try:
        slides = [
            {"page": i + 1, "text": _clean_text("\n".join(_shape_texts(slide.shapes)))}
            for i, slide in enumerate(prs.slides)
        ]
    except Exception as e:
        raise PptxParseError(f"슬라이드 텍스트 추출 실패: {e}") from e

    if not slides or all(not s["text"] for s in slides):
        raise UnprocessablePptxError("텍스트 없음(이미지 전용 덱 추정)")
    return slides


def _shape_texts(shapes) -> list[str]:
    """도형 목록에서 텍스트 조각을 모은다. 그룹은 재귀, 표는 행 단위로 펼친다."""
    parts: list[str] = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            parts.extend(_shape_texts(shape.shapes))
        elif shape.has_text_frame:
            parts.append(shape.text_frame.text)
        elif shape.has_table:
            for row in shape.table.rows:
                parts.append("\t".join(cell.text for cell in row.cells))
    return parts


def _clean_text(raw: str) -> str:
    """줄 단위 공백 정리 + 연속 빈 줄 제거. 줄 순서는 보존."""
    lines = [line.strip() for line in raw.splitlines()]
    cleaned: list[str] = []
    for line in lines:
        if line or (cleaned and cleaned[-1]):
            cleaned.append(line)
    return "\n".join(cleaned).strip()
