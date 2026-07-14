"""자료(PDF·PPTX) 업로드 + 백그라운드 파싱 회귀 테스트 (작업 3-1·3-2).

실행:
    cd backend
    .\\.venv\\Scripts\\python.exe -m pytest tests/test_materials.py -v

TestClient는 BackgroundTasks를 응답 후 동기 실행하므로, POST가 반환된 뒤엔
파싱이 이미 끝나 있다 → materials.status를 바로 확인할 수 있다.
"""

import io

import fitz  # PyMuPDF — 테스트용 PDF 즉석 생성
import pytest
from fastapi.testclient import TestClient
from pptx import Presentation  # 테스트용 PPTX 즉석 생성
from pptx.util import Inches
from sqlalchemy import delete

from app.core import storage
from app.db.models import Material, Team, TeamMember, User
from app.db.session import SessionLocal
from app.main import app
from tests.conftest import mark_email_verified

client = TestClient(app)


def _pdf(page_texts: list[str]) -> bytes:
    """각 문자열이 한 페이지. 빈 문자열 = 텍스트 없는 페이지(스캔본 흉내).

    주의: PyMuPDF 기본 폰트(helv)는 한글 글리프가 없어 한글은 점(·)으로 렌더된다.
    파서 자체는 언어 무관하므로 테스트 텍스트는 영문으로 쓴다."""
    doc = fitz.open()
    for t in page_texts:
        page = doc.new_page()
        if t:
            page.insert_text((72, 72), t)
    data = doc.tobytes()
    doc.close()
    return data


PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _pptx(slide_texts: list[str]) -> bytes:
    """각 문자열이 한 슬라이드. 빈 문자열 = 텍스트 없는 슬라이드(이미지 덱 흉내)."""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for t in slide_texts:
        slide = prs.slides.add_slide(blank)
        if t:
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
            box.text_frame.text = t
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _mkuser(u: str) -> str:
    r = client.post("/api/v1/auth/signup", json={"name": u, "username": u,
                    "password": "matr-pass-123", "email": f"{u}@t.io"})
    mark_email_verified(u)  # 로그인 차단(403) 우회 — email-verification-plan 작업 6
    return r.json()["user"]["id"]


def _auth(u: str) -> dict:
    tok = client.post("/api/v1/auth/login", json={"username": u, "password": "matr-pass-123"},
                      headers={"X-Client-Platform": "ios"}).json()["access_token"]
    return {"Authorization": f"Bearer {tok}"}


@pytest.fixture()
def ctx():
    """owner(발표자)·member·outsider. owner가 세션 하나 보유. → (session_id, storage_key)."""
    ids = {r: _mkuser(f"matr_{r}") for r in ("owner", "member", "outsider")}
    tid = client.post("/api/v1/teams", json={"name": "자료팀"}, headers=_auth("matr_owner")).json()["id"]
    with SessionLocal() as db:
        db.add(TeamMember(team_id=tid, user_id=ids["member"]))
        db.commit()
    sid = client.post(f"/api/v1/teams/{tid}/sessions",
                      json={"name": "발표", "personas": ["egen"], "question_count": 3,
                            "time_limit_minutes": 10, "mode": "realtime"},
                      headers=_auth("matr_owner")).json()["id"]
    yield sid, ids, tid
    with SessionLocal() as db:
        db.execute(delete(Team).where(Team.id == tid))
        db.execute(delete(User).where(User.username.ilike("matr_%")))
        db.commit()
    storage.delete(storage.material_key(sid))
    storage.delete(storage.material_key(sid, "pptx"))


def _upload(sid, pdf_bytes, who="matr_owner", filename="deck.pdf", ctype="application/pdf"):
    return client.post(f"/api/v1/sessions/{sid}/material",
                       files={"file": (filename, pdf_bytes, ctype)}, headers=_auth(who))


def _material(sid) -> Material:
    with SessionLocal() as db:
        return db.get(Material, sid)


class TestUploadSuccess:
    def test_valid_pdf_parses_to_ready(self, ctx):
        sid, _, _ = ctx
        r = _upload(sid, _pdf(["Cover slide text", "Problem definition text"]))
        assert r.status_code == 202
        assert r.json()["status"] == "queued"
        m = _material(sid)  # BackgroundTask 이미 실행됨
        assert m.status == "ready"
        assert m.page_count == 2
        assert m.progress == 1.0
        assert m.slides[0]["page"] == 1 and "Cover" in m.slides[0]["text"]
        assert m.error_code is None

    def test_file_saved_to_storage(self, ctx):
        sid, _, _ = ctx
        _upload(sid, _pdf(["x"]))
        assert storage.exists(storage.material_key(sid))

    def test_reupload_overwrites(self, ctx):
        sid, _, _ = ctx
        _upload(sid, _pdf(["first upload"]))
        _upload(sid, _pdf(["second try", "page two", "page three"]))
        m = _material(sid)
        assert m.page_count == 3  # 두 번째 업로드 반영
        assert "second" in m.slides[0]["text"]


class TestUploadFailureParsing:
    def test_scan_pdf_fails_unprocessable_not_500(self, ctx):
        """텍스트 레이어 없는 PDF → 업로드는 202, 파싱은 failed UNPROCESSABLE_PDF."""
        sid, _, _ = ctx
        r = _upload(sid, _pdf(["", ""]))  # 빈 페이지 = 스캔본 흉내
        assert r.status_code == 202  # 업로드 자체는 성공
        m = _material(sid)
        assert m.status == "failed"
        assert m.error_code == "UNPROCESSABLE_PDF"

    def test_corrupt_bytes_fails_parse_error(self, ctx):
        sid, _, _ = ctx
        r = _upload(sid, b"this is definitely not a pdf")
        assert r.status_code == 202
        m = _material(sid)
        assert m.status == "failed"
        assert m.error_code == "PDF_PARSE_ERROR"  # retry 대상

    def test_page_extraction_error_wrapped_as_parse_error(self, monkeypatch):
        """fitz.open은 성공하지만 페이지 추출이 raw 예외를 던지는 PDF(부분 손상)
        → 파서가 문서화된 계약대로 PdfParseError로 변환한다."""
        from app.services.material import PdfParseError, parse_pdf_to_slides

        def boom(self, *a, **kw):
            raise RuntimeError("malformed page stream")

        monkeypatch.setattr(fitz.Page, "get_text", boom)
        with pytest.raises(PdfParseError):
            parse_pdf_to_slides(_pdf(["ok page"]))

    def test_unexpected_job_exception_fails_not_stuck_processing(self, ctx, monkeypatch):
        """파싱 잡에서 예상 못한 예외가 나도 processing에 갇히지 않고 failed로
        끝난다(retry는 failed만 받으므로 stuck processing엔 복구 경로가 없다)."""
        import app.api.routes.materials as mat
        sid, _, _ = ctx

        def boom(_data, _ext):
            raise RuntimeError("unexpected bug in parser")

        monkeypatch.setattr(mat, "parse_material_to_slides", boom)
        r = _upload(sid, _pdf(["ok page"]))
        assert r.status_code == 202
        m = _material(sid)
        assert m.status == "failed"  # processing에 갇히면 회귀
        assert m.error_code == "PDF_PARSE_ERROR"


class TestUploadValidation:
    def test_non_pdf_extension_415(self, ctx):
        sid, _, _ = ctx
        r = _upload(sid, b"hello", filename="notes.txt", ctype="text/plain")
        assert r.status_code == 415
        assert r.json()["error"]["code"] == "UNSUPPORTED_MEDIA"

    def test_empty_file_400(self, ctx):
        sid, _, _ = ctx
        r = _upload(sid, b"")
        assert r.status_code == 400
        assert r.json()["error"]["code"] == "EMPTY_FILE"

    def test_oversize_413(self, ctx, monkeypatch):
        """20MB 상한 초과 → 413. 상한을 낮춰 대용량 페이로드 없이 경계 로직만 검증."""
        import app.api.routes.materials as mat
        monkeypatch.setattr(mat, "_MAX_BYTES", 100)
        r = _upload(sid=ctx[0], pdf_bytes=b"%PDF" + b"0" * 200)
        assert r.status_code == 413
        assert r.json()["error"]["code"] == "FILE_TOO_LARGE"


class TestUploadHardening:
    """재검증(2차) — 실전 엣지 + 동시성."""

    def test_octet_stream_content_type_accepted(self, ctx):
        """실제 업로드는 content-type이 octet-stream인 경우가 많다 — .pdf면 허용."""
        sid, _, _ = ctx
        r = _upload(sid, _pdf(["octet stream"]), ctype="application/octet-stream")
        assert r.status_code == 202
        assert _material(sid).status == "ready"

    def test_over_50_pages_unprocessable(self, ctx):
        """파서 페이지 상한(50) 초과 → failed UNPROCESSABLE_PDF (파이프라인 관통)."""
        sid, _, _ = ctx
        _upload(sid, _pdf([f"page {i}" for i in range(51)]))
        m = _material(sid)
        assert m.status == "failed" and m.error_code == "UNPROCESSABLE_PDF"

    def test_failed_then_reupload_recovers(self, ctx):
        """손상본으로 failed 된 뒤 정상 재업로드하면 ready로 복구된다."""
        sid, _, _ = ctx
        _upload(sid, b"garbage not pdf")
        assert _material(sid).status == "failed"
        _upload(sid, _pdf(["recovered"]))
        assert _material(sid).status == "ready"

    def test_missing_file_field_422(self, ctx):
        sid, _, _ = ctx
        assert client.post(f"/api/v1/sessions/{sid}/material",
                           headers=_auth("matr_owner")).status_code == 422

    def test_concurrent_first_upload_no_pk_race(self, ctx):
        """같은 세션 첫 업로드가 동시에 와도 500 없이 둘 다 202 (세션 행 잠금).
        재검증에서 발견: 잠금 전엔 materials_pkey UniqueViolation → 500 발생."""
        from concurrent.futures import ThreadPoolExecutor
        sid, _, _ = ctx
        body = _pdf(["concurrent"])
        with ThreadPoolExecutor(max_workers=2) as ex:
            futures = [ex.submit(_upload, sid, body) for _ in range(2)]
            codes = sorted(f.result().status_code for f in futures)
        assert codes == [202, 202]
        assert _material(sid).status == "ready"


class TestUploadPptx:
    """PPTX 자료 지원 — PDF와 동일 파이프라인(202 → 백그라운드 파싱 → 폴링)."""

    def test_valid_pptx_parses_to_ready(self, ctx):
        sid, _, _ = ctx
        r = _upload(sid, _pptx(["Cover slide text", "Body slide text"]),
                    filename="deck.pptx", ctype=PPTX_MIME)
        assert r.status_code == 202
        assert r.json()["status"] == "queued"
        m = _material(sid)
        assert m.status == "ready"
        assert m.page_count == 2
        assert m.slides[0]["page"] == 1 and "Cover" in m.slides[0]["text"]
        assert m.storage_key.endswith("material.pptx")
        assert m.error_code is None

    def test_octet_stream_pptx_extension_accepted(self, ctx):
        """PDF와 동일 — content-type이 octet-stream이어도 .pptx 확장자면 허용."""
        sid, _, _ = ctx
        r = _upload(sid, _pptx(["octet stream"]), filename="deck.pptx",
                    ctype="application/octet-stream")
        assert r.status_code == 202
        assert _material(sid).status == "ready"

    def test_textless_deck_fails_unprocessable(self, ctx):
        """텍스트 없는 덱(이미지 전용) → 업로드는 202, 파싱은 failed UNPROCESSABLE_PPTX."""
        sid, _, _ = ctx
        r = _upload(sid, _pptx(["", ""]), filename="deck.pptx", ctype=PPTX_MIME)
        assert r.status_code == 202
        m = _material(sid)
        assert m.status == "failed"
        assert m.error_code == "UNPROCESSABLE_PPTX"

    def test_corrupt_pptx_parse_error_and_retryable(self, ctx):
        sid, _, _ = ctx
        _upload(sid, b"this is definitely not a pptx", filename="deck.pptx", ctype=PPTX_MIME)
        m = _material(sid)
        assert m.status == "failed"
        assert m.error_code == "PPTX_PARSE_ERROR"  # retry 대상
        r = client.post(f"/api/v1/sessions/{sid}/material/retry", headers=_auth("matr_owner"))
        assert r.status_code == 202  # pptx도 재시도 경로가 열려 있다

    def test_over_50_slides_unprocessable(self, ctx):
        """슬라이드 상한(50) 초과 → failed UNPROCESSABLE_PPTX (파이프라인 관통)."""
        sid, _, _ = ctx
        _upload(sid, _pptx([f"slide {i}" for i in range(51)]),
                filename="deck.pptx", ctype=PPTX_MIME)
        m = _material(sid)
        assert m.status == "failed" and m.error_code == "UNPROCESSABLE_PPTX"

    def test_legacy_ppt_415(self, ctx):
        """레거시 .ppt(바이너리)는 미지원 — 415로 즉시 거부."""
        sid, _, _ = ctx
        r = _upload(sid, b"\xd0\xcf\x11\xe0 legacy ppt binary", filename="deck.ppt",
                    ctype="application/vnd.ms-powerpoint")
        assert r.status_code == 415
        assert r.json()["error"]["code"] == "UNSUPPORTED_MEDIA"

    def test_cross_format_reupload_no_orphan_file(self, ctx):
        """pdf → pptx 재업로드 시 이전 material.pdf가 스토리지에 고아로 남지 않는다."""
        sid, _, _ = ctx
        _upload(sid, _pdf(["pdf first"]))
        assert storage.exists(storage.material_key(sid, "pdf"))
        _upload(sid, _pptx(["pptx second"]), filename="deck.pptx", ctype=PPTX_MIME)
        assert not storage.exists(storage.material_key(sid, "pdf"))  # 고아 방지
        assert storage.exists(storage.material_key(sid, "pptx"))
        m = _material(sid)
        assert m.status == "ready" and m.storage_key.endswith(".pptx")


class TestUploadPermission:
    def test_member_not_owner_403(self, ctx):
        sid, _, _ = ctx
        assert _upload(sid, _pdf(["x"]), who="matr_member").status_code == 403

    def test_outsider_404(self, ctx):
        sid, _, _ = ctx
        assert _upload(sid, _pdf(["x"]), who="matr_outsider").status_code == 404

    def test_requires_auth(self, ctx):
        sid, _, _ = ctx
        r = client.post(f"/api/v1/sessions/{sid}/material",
                        files={"file": ("d.pdf", _pdf(["x"]), "application/pdf")})
        assert r.status_code == 401
